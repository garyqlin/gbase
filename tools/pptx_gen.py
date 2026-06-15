# SPDX-License-Identifier: MIT
"""PPT 演示文稿生成工具"""

import json
import logging
import os
import subprocess

from lib.toolkit import tool

logger = logging.getLogger(__name__)


@tool()
async def gen_pptx(title: str = "演示文稿", slides: list = None, output_path: str = "") -> dict:
    """
    生成 PowerPoint (.pptx) 文件。

    Args:
        title: 文件名（不含路径）
        slides: 幻灯片列表，每项为 {"title": "封面", "content": ["第一行","第二行"], "layout": "title|content|two"}
                layout: title=标题页, content=内容页(默认), two=两栏, blank=空白
        output_path: 输出路径

    Returns:
        {"path": "...", "size": 12345, "slides": 3}
    """
    if not slides:
        slides = [{"title": title, "content": ["欢迎"], "layout": "title"}]

    if not output_path:
        output_path = os.path.expanduser(f"~/Downloads/{title}.pptx")

    script = f"""import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides = {json.dumps(slides)}

for sdata in slides:
    layout_name = sdata.get("layout", "content")
    if layout_name == "title":
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        # 标题居中大号
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sdata.get("title", "")
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
        p.alignment = PP_ALIGN.CENTER

        content = sdata.get("content", [])
        if content:
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(2))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for ci, line in enumerate(content):
                if ci == 0:
                    p2 = tf2.paragraphs[0]
                else:
                    p2 = tf2.add_paragraph()
                p2.text = str(line)
                p2.font.size = Pt(24)
                p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                p2.alignment = PP_ALIGN.CENTER
    else:
        # 内容页
        if layout_name == "blank":
            layout = prs.slide_layouts[6]
        else:
            layout = prs.slide_layouts[1]  # title and content
        slide = prs.slides.add_slide(layout)

        # 标题
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = sdata.get("title", "")

        content = sdata.get("content", [])
        if content:
            body = slide.placeholders[1]  # content placeholder
            if body:
                tf = body.text_frame
                tf.clear()
                for ci, line in enumerate(content):
                    if ci == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = str(line)
                    p.font.size = Pt(18)
                    p.space_after = Pt(8)

prs.save({json.dumps(output_path)})
print("OK:" + {json.dumps(output_path)})
"""

    result = subprocess.run(["python3", "-c", script], capture_output=True, text=True, timeout=30)

    if result.returncode != 0:
        return {"error": result.stderr.strip() or result.stdout.strip()}

    if os.path.isfile(output_path):
        size = os.path.getsize(output_path)
        return {"path": output_path, "size": size, "slides": len(slides)}
    return {"error": "文件未生成"}
